# Seasonal trends
# Weekly conference
#

import os
import pandas as pd
# import qgis
import psycopg2
import datetime as dt
import numpy as np
import sqlalchemy
import matplotlib.pyplot as plt
from PIL import Image
from pptx import Presentation
import seaborn as sns
from rasterstats import zonal_stats
import geopandas as gpd
import csv

def ppt_generator(path, years, range, df_sql, nat2k_year, nat2k_week):
	os.chdir(path)
	vc_day = (dt.datetime.today() + dt.timedelta(1)).strftime('%d.%m.%Y')
	vc_yesterday = dt.datetime.today().strftime('%d.%m.%Y')
	vc_7daysago = (dt.datetime.today() - dt.timedelta(7)).strftime('%d.%m.%Y')
	prs = Presentation('vc_template.pptx')

	slide3, slide4 = weekly_sum_count(years, range, path, df_sql)

	# slide 1
	title_slide_layout = prs.slide_layouts[0]
	slide = prs.slides.add_slide(title_slide_layout)
	title = slide.shapes.title
	subtitle = slide.placeholders[1]
	title.text = "Videoconference - {}".format(vc_day)
	title.width
	subtitle.text = "https://effis.jrc.ec.europa.eu" + "\n" + "jrc-effis@ec.europa.eu"

	# slide 2
	slide = prs.slides.add_slide(title_slide_layout)
	title = slide.shapes.title
	title.text = "EFFIS" + "\n" + "Burnt areas mapped until {}".format(vc_yesterday)

	# slide 3
	slide_layout = prs.slide_layouts[1]
	slide = prs.slides.add_slide(slide_layout)
	title = slide.shapes.title
	title.text = "\n\n" + "Burnt area mapped in EFFIS (fires larger than 30 ha)" + "\n" + "from 01.01.2020 to {}".format(vc_yesterday) + "" \
				 "\n\n\n\n" + "Total Area (ha) European Union: {:,}".format(slide3['AREA_HA'][0]) + "\n\n" + "" \
				"Total Area (ha) Other European countries: {:,}".format(slide3['AREA_HA'][1]) + "\n\n" + "" \
				"Total Area (ha) Middle East and North Africa: {:,}".format(slide3['AREA_HA'][2])
	subtitle = slide.placeholders[1]
	subtitle.text = 'Total Area Burnt (ha) mapped in EU Natura2000 sites: {:,} ({}%)'.format(int(nat2k_year['tot_nat2k_ha'][0]), int((nat2k_year['tot_nat2k_ha'][0]/slide3['AREA_HA'][0])*100))

	# slide 4
	slide_layout = prs.slide_layouts[1]
	slide = prs.slides.add_slide(slide_layout)
	title = slide.shapes.title
	# check in case during the last week there was no burnt areas in EU, EU_non, ME_AF
	try:
		slide4['AREA_HA'][0]
	except:
		slide4 = slide4.append(pd.DataFrame.from_dict({'EU': [0]}, orient='index', columns=['AREA_HA']))
	try:
		slide4['AREA_HA'][1]
	except:
		slide4 = slide4.append(pd.DataFrame.from_dict({'EU_non': [0]}, orient='index', columns=['AREA_HA']))
	try:
		slide4['AREA_HA'][2]
	except:
		slide4 = slide4.append(pd.DataFrame.from_dict({'ME_AF': [0]}, orient='index', columns=['AREA_HA']))
	title.text = "\n\n" + "Burnt area mapped in EFFIS (fires larger than 30 ha)" + "\n" + "from {} to {}".format(
		vc_7daysago, vc_yesterday) + "" \
		"\n\n\n\n" + "Total Area (ha) European Union: {:,}".format(slide4['AREA_HA'][0]) + "\n\n" + "" \
		"Total Area (ha) Other European countries: {:,}".format(
		slide4['AREA_HA'][1]) + "\n\n" + "" \
		"Total Area (ha) Middle East and North Africa: {:,}".format(
		slide4['AREA_HA'][2])
	subtitle = slide.placeholders[1]
	if nat2k_week['tot_nat2k_ha'][0] == None:
		nat2k_week['tot_nat2k_ha'][0] = 0
	try:
		percentage = int((nat2k_week['tot_nat2k_ha'][0] / slide4['AREA_HA'][0]) * 100)
	except:
		percentage = 0
	subtitle.text = 'Total Area Burnt (ha) mapped in EU Natura2000 sites: {:,} ({}%)'.format(int(nat2k_week['tot_nat2k_ha'][0]), percentage)

	# slide 5
	slide_layout = prs.slide_layouts[2]
	slide = prs.slides.add_slide(slide_layout)
	title = slide.shapes.title
	title.text = "Weekly evolution of the burned area mapped in EFFIS \n (EU countries)"
	picture_placeholder = slide.placeholders[1]
	picture_placeholder
	picture_placeholder.shape_type
	placeholder_picture = picture_placeholder.insert_picture('slide5.png')

	# slide 6
	slide_layout = prs.slide_layouts[2]
	slide = prs.slides.add_slide(slide_layout)
	title = slide.shapes.title
	title.text = "Weekly evolution of the burned area mapped in EFFIS \n (non EU countries)"
	picture_placeholder = slide.placeholders[1]
	picture_placeholder
	picture_placeholder.shape_type
	placeholder_picture = picture_placeholder.insert_picture('slide6.png')

	# slide 7
	slide_layout = prs.slide_layouts[2]
	slide = prs.slides.add_slide(slide_layout)
	title = slide.shapes.title
	title.text = "Cumulative Burnt Area mapped in EFFIS \n (fires larger than 30 ha) - EU Countries"
	picture_placeholder = slide.placeholders[1]
	picture_placeholder
	picture_placeholder.shape_type
	placeholder_picture = picture_placeholder.insert_picture('slide7.png')

	# slide 8
	slide_layout = prs.slide_layouts[2]
	slide = prs.slides.add_slide(slide_layout)
	title = slide.shapes.title
	title.text = "Cumulative number of fires mapped in EFFIS \n (fires larger than 30 ha) - EU Countries"
	picture_placeholder = slide.placeholders[1]
	picture_placeholder
	picture_placeholder.shape_type
	placeholder_picture = picture_placeholder.insert_picture('slide8.png')

	# slide 9
	slide_layout = prs.slide_layouts[2]
	slide = prs.slides.add_slide(slide_layout)
	title = slide.shapes.title
	title.text = "Active fires detected in EFFIS – last 7 days"
	# picture_placeholder = slide.placeholders[1]
	# picture_placeholder
	# picture_placeholder.shape_type
	# placeholder_picture = picture_placeholder.insert_picture('slide8.png')

	# slide 10
	slide_layout = prs.slide_layouts[2]
	slide = prs.slides.add_slide(slide_layout)
	title = slide.shapes.title
	title.text = "Burnt Areas mapped in EFFIS – last 7 days"

	# slide 11
	slide_layout = prs.slide_layouts[2]
	slide = prs.slides.add_slide(slide_layout)
	title = slide.shapes.title
	title.text = "Burnt Areas mapped in EFFIS – Since Jan. 1, 2020"

	prs.save('vc_{}.pptx'.format(vc_day))

def db_connection():
	try:
		connection = psycopg2.connect(user = "pieralberto",
									  password = "4q2WK4MZMs",
									  host = "db1.wild-fire.eu",
									  port = "5432",
									  database = "e1gwis")

		cursor = connection.cursor()
		# Print PostgreSQL Connection properties
		print(connection.get_dsn_parameters(), "\n")

		# Print PostgreSQL version
		cursor.execute("SELECT version();")
		record = cursor.fetchone()
		print("You are connected to - ", record, "\n")
		cursor.execute("select relname from pg_class where relkind='r' and relname !~ '^(pg_|sql_)';")
		# print(cursor.fetchall())
		# temp = cursor.fetchall()

		engine_string = "postgresql+psycopg2://{user}:{password}@{host}:{port}/{database}".format(
			user = "pieralberto",
			password = "4q2WK4MZMs",
			host = "db1.wild-fire.eu",
			port = "5432",
			database = "e1gwis",
		)
		# print('creating sqlalchemy engine')
		engine = sqlalchemy.create_engine(engine_string)
		# read a table from database into pandas dataframe, replace "tablename" with your table name
		global df_sql
		df_sql = pd.read_sql_query('select * from gw_burntarea_effis.ba_oracle_export_year', con = engine)

		nat2k_year = pd.read_sql_query('select sum(clc.natura2k*ba.area_ha) as '
									   'tot_nat2k_ha from gw_burntarea_effis.ba_final as ba, '
									   'gw_burntarea_effis.ba_stats_clc as clc where ba.id=clc.ba_id and '
									   'ba.initialdate>=\'2020-01-01\' and clc.natura2k is not null;', con = engine)
		nat2kweek = pd.read_sql_query('select sum(clc.natura2k*ba.area_ha) as tot_nat2k_ha from '
									  'gw_burntarea_effis.ba_final as ba, gw_burntarea_effis.ba_stats_clc '
									  'as clc where ba.id=clc.ba_id and ba.initialdate>=(now()-interval \'7 days\')::date '
									  'and clc.natura2k is not null;', con = engine)
		print("df_sql successfully downloaded")
		global gdf_sql
		gdf_sql = gpd.read_postgis('select * from gw_burntarea_effis.ba_oracle_export_year', con=engine)
		print("gdf_sql successfully downloaded")
		global tab_nations  # Table associating the
		tab_nations = pd.read_sql_query('select * from "burnt_areas"."Tab_Elenco_Nazioni"', con=engine)
		print("tab_nations successfully downloaded")
	except (Exception, psycopg2.Error) as error:
		print("Error while connecting to PostgreSQL", error)
	finally:
		# closing database connection.
		if (connection):
			cursor.close()
			connection.close()
			print("PostgreSQL connection now closed")
	return df_sql, nat2k_year, nat2kweek, gdf_sql, tab_nations

def dates_range_set(years):
	d = {}
	for i in years:
		first_day = dt.datetime.strptime("01-01-{}".format(i), "%d-%m-%Y")
		dates = []
		for j in (np.arange(53)):
			days = 7 * j
			day = first_day + dt.timedelta(int(days))
			dates.append(day)
		d["range_dates_{}".format(i)] = dates
	return d

def get_text_positions(x_data, y_data, txt_width, txt_height):
	a = list(zip(y_data, x_data))
	text_positions = y_data.copy()
	for index, (y, x) in enumerate(a):
		local_text_positions = [i for i in a if i[0] > (y - txt_height)
							and (abs(i[1] - x) < txt_width * 3.5) and i != (y,x)]
		if local_text_positions:
			sorted_ltp = sorted(local_text_positions)
			if abs(sorted_ltp[0][0] - y) < txt_height: #True == collision
				differ = np.diff(sorted_ltp, axis=0)
				a[index] = (sorted_ltp[-1][0] + txt_height, a[index][1])
				text_positions[index] = sorted_ltp[-1][0] + txt_height
				for k, (j, m) in enumerate(differ):
					#j is the vertical distance between words
					if j > txt_height * 2.5: #if True then room to fit a word in
						a[index] = (sorted_ltp[k][0] + txt_height, a[index][1])
						text_positions[index] = sorted_ltp[k][0] + txt_height
						break
	return text_positions

def text_plotter(x_data, y_data, text_positions, axis,txt_width,txt_height):
	for x,y,t in zip(x_data, y_data, text_positions):
		axis.text(x - .03, 1.02*t, '%d'%int(y),rotation=0, color='black', fontsize=13)
		if y != t:
			axis.arrow(x, t+20,0,y-t, color='blue',alpha=0.6, width=txt_width*0.0,
					   head_width=.02, head_length=txt_height*0.5,
					   zorder=0,length_includes_head=True)

def labeled_graphs(path, df_cumulative, title, png_name):
	print('Plotting marked and labeled graphs')
	font = {'family': 'serif',
			'color': 'black',
			'weight': 'normal',
			'size': 20,
			}
	df_cumulative = df_cumulative.truncate(after=(dt.datetime.now() - dt.timedelta(days=0)).strftime("%Y-%m-%d"))
	# print("\n Df cumulative slide 5 truncate")
	# print(df_cumulative)
	# Original labeling
	# df_cumulative['Area_HA'].plot(legend=False, figsize=[10,7], color='r', marker='o') # figsize=[8,4],
	# plt.title('2020 burnt area EU countries', fontdict=font)
	# plt.grid(color='grey', linestyle='-', linewidth=1)
	# plt.ylabel('Area (ha)', fontdict=font)
	# plt.xlabel('Time', fontdict=font)
	# area_labels = df_cumulative['Area_HA'].astype(int).apply(lambda x: "{:,}".format(x))
	# for a, b, c in zip(df_cumulative.index, df_cumulative['Area_HA'], area_labels):
	# 	label = "{}".format(c)
	# 	plt.annotate(label,  # this is the text
	# 				 (a, b),  # this is the point to label
	# 				 textcoords="offset points",  # how to position the text
	# 				 xytext=(0, 10),  # distance from text to points (x,y)
	# 				 ha='center')  # horizontal alignment can be left, right or center
	# Adaptive labeling
	fig = plt.figure(figsize=[14, 9])
	ax = fig.add_subplot(111)
	x_data = np.arange(len(df_cumulative.index))
	y_data = df_cumulative['Area_HA']
	ax.plot(x_data, y_data, color='red', marker='o');
	ax.grid(color='grey', linestyle='-', linewidth=1, alpha=0.2)
	plt.title(title, fontdict=font)
	plt.ylabel('Area (ha)', fontdict=font)
	plt.xlabel('Time', fontdict=font)
	# ticks
	dates_label = df_cumulative.index.strftime("%d %b")
	plt.xticks(ticks=x_data, labels=dates_label, rotation='vertical')
	# labels on markers
	txt_height = 0.04 * (plt.ylim()[1] - plt.ylim()[0])
	txt_width = 0.02 * (plt.xlim()[1] - plt.xlim()[0])
	marker_labels = df_cumulative['Area_HA'].astype(int).apply(lambda x: "{:,}".format(x))
	text_positions = get_text_positions(x_data, y_data, txt_width, txt_height)
	text_plotter(x_data, y_data, text_positions, plt, txt_width, txt_height)
	# end of labeling
	plt.savefig(png_name)
	im = Image.open(path + png_name)
	# im.show()
	plt.clf()

def unlabeled_graphs(path, df_comparison, title, png_name, ylabel):
	print('Plotting unlabeled graphs')
	plt.clf()
	ylabel = ylabel
	font = {'family': 'serif',
			'color': 'black',
			'weight': 'normal',
			'size': 20,
			}
	df_comparison = df_comparison.rename(
		columns={"Area_HA": "2020 burnt area", "hist_average_from_2008": "2008-2019 average burnt area"})
	df_comparison = df_comparison.truncate(after=(dt.datetime.now() - dt.timedelta(days=1)).strftime("%Y-%m-%d"))
	# print(df_comparison)
	df_comparison['2020 burnt area'].plot(legend=True, figsize=[10, 7], color='r')
	df_comparison['2008-2019 average burnt area'].plot(legend=True, figsize=[10, 7], color='b')
	plt.ylabel(ylabel, fontdict=font)
	plt.xlabel('Time', fontdict=font)
	plt.title(title, fontdict=font)
	plt.grid(color='grey', linestyle='-', linewidth=1)
	plt.savefig(png_name)
	from PIL import Image
	im = Image.open(path + png_name)
	# im.show()
	plt.clf()

def weekly_sum_count(years, range, path, df_sql):
	os.chdir(path)
	year = str(dt.datetime.now().year)
	df = df_sql
	# print(df.columns)
	df['AREA_HA'] = df['AREA_HA'].astype(int)
	df_sql.to_csv('df_sql.csv')
	df_elenco_nazioni = pd.read_csv('tab_elenco_nazioni.csv')
	# Conversions from string to date format
	df['FIREDATE'] = pd.to_datetime(df['FIREDATE'])
	df['lastfiredate'] = pd.to_datetime(df['lastfiredate'])
	df['LASTUPDATE'] = pd.to_datetime(df['LASTUPDATE'])

	list_eunon = ['EU','EU_non','ME_AF']
	print('The dataset will be divided in three macroregions: {}'.format(list_eunon))
	# dataframes definition
	df = df.merge(df_elenco_nazioni, how = 'inner', left_on = 'COUNTRY', right_on = 'NUTS0_CODE')
	print('\n --- Count of burnt areas --- \n')
	print(df.groupby([pd.Grouper(key = 'EU_nonEU')])['AREA_HA'].count())
	df_temp = pd.DataFrame(df.groupby([pd.Grouper(key = 'EU_nonEU')])['AREA_HA'].count())
	# print('Sum of count: \n')
	# print(df_temp.append(df_temp.sum(numeric_only=True), ignore_index=True))
	# print('\n')
	df = df.loc[df['AREA_HA'] >= 30] # areas over 30ha to harmonize the comparison with historical data
	df_eu = df.loc[df['EU_nonEU'] == 'EU']
	df_eunon = df.loc[df['EU_nonEU'] == 'EU_non']
	df_meaf = df.loc[df['EU_nonEU'] == 'ME_AF']
	df_eunon_meaf = pd.concat([df_eunon, df_meaf])
	df_history = pd.read_csv('ba_2008_2019.csv')
	df_history = df_history.loc[df_history['Area_HA'] >= 30]  # areas over 30ha to harmonize the comparison with historical data
	df_history['FireDate'] = pd.to_datetime(df_history['FireDate'])
	df_history = df_history.merge(df_elenco_nazioni, how = 'inner', left_on = 'Country', right_on = 'NUTS0_CODE')

	#slide 3
	df_slide3 = pd.DataFrame(df.groupby([pd.Grouper(key='EU_nonEU')])['AREA_HA'].sum())
	# print(df_slide3)
	pd.DataFrame(df.groupby([pd.Grouper(key = 'EU_nonEU')])['AREA_HA'].sum()).to_csv('slide3.csv')

	#slide 4
	today = dt.datetime.today()
	seven_days_ago = today - dt.timedelta(days = 8) # 8 days because the code is run during the afternoon of the day before the videoconference
	df_last_week = df.loc[df['FIREDATE'] >= seven_days_ago]
	df_last_week.to_csv('delta.csv')
	df_slide4 = pd.DataFrame(df_last_week.groupby([pd.Grouper(key='EU_nonEU')])['AREA_HA'].sum())
	pd.DataFrame(df_last_week.groupby([pd.Grouper(key = 'EU_nonEU')])['AREA_HA'].sum()).to_csv('slide4.csv')

	#slide 5
	range_2020 = pd.DataFrame(range.get('range_dates_2020'))
	range_2020.columns = ['bins']
	dates_tag = []
	for i in range_2020['bins']:
		k = str(i)[0:10]
		dates_tag.append(dt.datetime.strptime(k, '%Y-%m-%d'))
	# print(dates_tag)

	df_area = pd.DataFrame(df_eu.groupby(pd.cut(df_eu['FIREDATE'], range_2020['bins'], right = 0, include_lowest = 1))['AREA_HA'].sum())
	df_count = pd.DataFrame(df_eu.groupby(pd.cut(df_eu['FIREDATE'], range_2020['bins'], right = 0, include_lowest = 1))['AREA_HA'].count())
	# for aggregation on weekdays uncomment the next 2 lines
	# df_area = pd.DataFrame(df_eu.groupby([pd.Grouper(key = 'FIREDATE', freq = 'W-MON')])['AREA_HA'].sum().cumsum())
	# df_count = pd.DataFrame(df_eu.groupby([pd.Grouper(key = 'FIREDATE', freq = 'W-MON')])['AREA_HA'].count().cumsum())
	df_cumulative = df_area.merge(df_count, how = 'inner', on = 'FIREDATE')
	df_cumulative.index.name = 'Period'
	df_cumulative.columns = ['Area_HA', 'Count']
	df_cumulative.to_csv('slide5.csv'.format(year))
	dat = pd.DataFrame(dates_tag)
	dat.columns = ['Dates_tag']
	df_cumulative.insert(2, "Single dates", dates_tag[1:], True)
	df_cumulative = df_cumulative.reset_index()
	df_cumulative = df_cumulative.set_index('Single dates')
	# print('\n Cumulative slide 5')
	# print(df_cumulative)
	# plotting
	print('Now plotting slide 5')
	# print(df_cumulative)
	title = '2020 burnt area in EU countries'
	png_name = 'slide5.png'
	labeled_graphs(path, df_cumulative, title, png_name)
	# font = {'family': 'serif',
	# 		'color': 'black',
	# 		'weight': 'normal',
	# 		'size': 20,
	# 		}

	# slide 6
	df_area = pd.DataFrame(df_eunon_meaf.groupby(pd.cut(df_eunon_meaf['FIREDATE'], range_2020['bins'], right = 0, include_lowest = 1))['AREA_HA'].sum())
	df_count = pd.DataFrame(df_eunon_meaf.groupby(pd.cut(df_eunon_meaf['FIREDATE'], range_2020['bins'], right = 0, include_lowest = 1))['AREA_HA'].count())
	# for aggregation on weekdays uncomment the next 2 lines
	# df_area = pd.DataFrame(df_eu.groupby([pd.Grouper(key = 'FIREDATE', freq = 'W-MON')])['AREA_HA'].sum().cumsum())
	# df_count = pd.DataFrame(df_eu.groupby([pd.Grouper(key = 'FIREDATE', freq = 'W-MON')])['AREA_HA'].count().cumsum())
	df_cumulative = df_area.merge(df_count, how = 'inner', on = 'FIREDATE')
	df_cumulative.index.name = 'Period'
	df_cumulative.columns = ['Area_HA', 'Count']
	df_cumulative.to_csv('slide6.csv')
	df_cumulative.insert(2, "Single dates", dates_tag[1:], True)
	df_cumulative = df_cumulative.reset_index()
	df_cumulative = df_cumulative.set_index('Single dates')
	# plotting
	print('Now plotting slide 6')
	plt.clf()
	title = '2020 burnt area in non EU countries'
	png_name = 'slide6.png'
	# print(df_cumulative)
	labeled_graphs(path, df_cumulative, title, png_name)

	#slide 7
	df_history_eu = df_history.loc[df_history['EU_nonEU'] == 'EU']
	df_area_2020 = pd.DataFrame(
		df_eu.groupby(pd.cut(df_eu['FIREDATE'], range_2020['bins'], right = 0, include_lowest = 1))['AREA_HA'].sum().cumsum())
	df_area_2020.index.name = 'Period'
	df_area_2020.columns = ['Area_HA']
	hist_stats = {}
	for year in years:
		range_dates = pd.DataFrame(range.get('range_dates_{}'.format(year)))
		range_dates.columns = ['bins']
		df_year_temp = df_history_eu.loc[df_history_eu['YearSeason'] == year]
		df_year = df_year_temp.groupby(pd.cut(df_year_temp['FireDate'], range_dates['bins'], right = 0, include_lowest = 1))['Area_HA'].sum().cumsum()
		df_year.index.name = 'Period'
		df_year.columns = ['Area_HA']
		hist_stats["{}".format(year)] = df_year
	del hist_stats['2020']
	df_hist = df_area_2020.copy()
	for i in hist_stats.keys():
		df_hist[str(i)] = hist_stats.get(i).values

	# seaborn graph
	plt.clf()
	df_sea = df_hist
	col_2020 = df_sea['Area_HA']
	cols = list(df_sea.columns.values)
	cols.pop(cols.index('Area_HA'))
	df_sea = df_sea[cols + ['Area_HA']]
	df_sea = df_sea.rename(columns={"Area_HA": "2020"})
	df_sea.insert(0, "Week", dates_tag[1:], True)
	df_sea = df_sea.reset_index()
	df_sea = df_sea.drop('Period', axis=1)
	df_sea['Week'] = df_sea['Week'].dt.strftime('%m-%d')
	df_sea = df_sea.set_index('Week')
	# print(df_sea)
	df_sea.to_csv('sea.csv')
	midpoint = (df_sea.values.max() - df_sea.values.min()) / 4
	# print('midpoint is {}'.format(midpoint))
	seafig = sns.heatmap(df_sea, annot=True, fmt='g', center = midpoint, cmap='coolwarm').get_figure()
	seafig.savefig('sea.png')
	# end of seaborn graph

	df_hist = df_hist.drop('Area_HA', axis=1)
	df_hist['hist_average_from_2008'] = df_hist.mean(numeric_only = True, axis = 1).round(0)
	df_hist['hist_average_from_2008'] = df_hist['hist_average_from_2008'].astype(int)
	df_area_2020 = df_area_2020.merge(df_hist['hist_average_from_2008'], how = 'inner', left_on = 'Period', right_on = 'Period')
	df_area_2020.to_csv('slide7.csv')
	# plotting
	df_area_2020.insert(2, "Single dates", dates_tag[:-1], True)
	df_area_2020 = df_area_2020.reset_index()
	df_area_2020 = df_area_2020.set_index('Single dates')
	print('Now plotting slide 7')
	plt.clf()
	df_area_2020 = df_area_2020.rename(columns={"Area_HA": "2020 burnt area", "hist_average_from_2008": "2008-2019 average burnt area"})
	df_area_2020 = df_area_2020.truncate(after=(dt.datetime.now() - dt.timedelta(days = 1)).strftime("%Y-%m-%d"))
	# print(df_area_2020)
	title = 'Cumulative burnt area in EU countries'
	png_name = 'slide7.png'
	ylabel = 'Area (ha)'
	unlabeled_graphs(path, df_area_2020, title, png_name, ylabel)

	# slide 8
	df_count_2020 = pd.DataFrame(
		df_eu.groupby(pd.cut(df_eu['FIREDATE'], range_2020['bins'], right = 0, include_lowest = 1))['AREA_HA'].count().cumsum())
	df_count_2020.index.name = 'Period'
	df_count_2020.columns = ['count']
	hist_stats = {}
	for year in years:
		range_dates = pd.DataFrame(range.get('range_dates_{}'.format(year)))
		range_dates.columns = ['bins']
		df_year_temp = df_history_eu.loc[df_history_eu['YearSeason'] == year]
		df_year = df_year_temp.groupby(pd.cut(df_year_temp['FireDate'], range_dates['bins'], right = 0, include_lowest = 1))['Area_HA'].count().cumsum()
		df_year.index.name = 'Period'
		df_year.columns = ['count']
		hist_stats["{}".format(year)] = df_year
	del hist_stats['2020']
	df_hist = df_count_2020.copy()
	for i in hist_stats.keys():
		df_hist[str(i)] = hist_stats.get(i).values
	df_hist = df_hist.drop('count', axis = 1)
	df_hist['hist_average_from_2008'] = df_hist.mean(numeric_only = True, axis = 1).round(0)
	df_hist['hist_average_from_2008'] = df_hist['hist_average_from_2008'].astype(int)
	df_count_2020 = df_count_2020.merge(df_hist['hist_average_from_2008'], how = 'inner', left_on = 'Period',
										right_on = 'Period')
	df_count_2020.to_csv('slide8.csv')
	# plotting
	df_count_2020.insert(2, "Single dates", dates_tag[:-1], True)
	df_count_2020 = df_count_2020.reset_index()
	df_count_2020 = df_count_2020.set_index('Single dates')
	print('Now plotting slide 8')
	plt.clf()
	df_count_2020 = df_count_2020.rename(
		columns={"count": "2020 burnt area", "hist_average_from_2008": "2008-2019 average burnt area"})
	df_count_2020 = df_count_2020.truncate(after=(dt.datetime.now() - dt.timedelta(days = 1)).strftime("%Y-%m-%d"))
	# print(df_count_2020)
	title = 'Cumulative number of fires in EU countries'
	png_name = 'slide8.png'
	ylabel = 'Fires count'
	unlabeled_graphs(path, df_count_2020, title, png_name, ylabel)
	#
	# df_count_2020['2020 burnt area'].plot(legend=True, figsize=[10, 7], color='r')
	# df_count_2020['2008-2019 average burnt area'].plot(legend=True, figsize=[10, 7], color='b')
	# plt.ylabel('Fires count', fontdict=font)
	# plt.xlabel('Time', fontdict=font)
	# plt.title('Cumulative number of fires in EU countries', fontdict=font)
	# plt.grid(color='grey', linestyle='-', linewidth=1)
	# plt.savefig('slide8.png')
	# from PIL import Image
	# im = Image.open(path+"slide8.png")
	# im.show()

	return df_slide3, df_slide4

def corine_stats(burnt_areas, corine, dict_path, tab_nations):
	ba = burnt_areas
	ba = ba.merge(tab_nations, how='inner', left_on='COUNTRY', right_on='NUTS0_CODE')
	# print('--- Coulmns of burnt areas ---')
	# print(ba.columns)
	# ba = gpd.read_file(burnt_areas)
	ba = ba.sort_values('id')
	ba_eu = ba.loc[ba['EU_nonEU'] == 'EU']
	ba_eu_non = ba.loc[ba['EU_nonEU'] == 'EU_non']
	ba_mena = ba.loc[(ba['EU_nonEU'] != 'EU') & (ba['EU_nonEU'] != 'EU_non')]
	os.chdir(dict_path)
	f = csv.reader(open('corine_key.csv'))
	corine_dict = {}
	for row in f:
		corine_dict[row[0]] = row[1]
	corine_dict['1'] = corine_dict.pop('ï»¿1')
	corine_dict = {int(k): v for k, v in corine_dict.items()}
	ba_group = [ba_eu, ba_eu_non, ba_mena]
	ba_names = ['ba_eu', 'ba_eu_non', 'ba_mena']
	for i,j in zip(ba_group,ba_names):
		# ba_stats = zonal_stats(i, corine, categorical=True, category_map=corine_dict)
		ba_raw = zonal_stats(i, corine, categorical=True)
		df_ba_raw = pd.DataFrame.from_dict(ba_raw)
		df_ba_raw = df_ba_raw.rename(columns=corine_dict)
		df_ba_raw = df_ba_raw.T
		df_mapped = df_ba_raw.groupby(df_ba_raw.index).sum()
		df_mapped = df_mapped.T
		df = df_mapped
		if 0 in df.columns:
			df = df.rename(columns={0: 'Unclassified'})
		df = df.append(df.sum(numeric_only=True), ignore_index=True)
		pd.options.display.float_format = '{:,.0f}'.format
		df.to_csv('corine_{}.csv'.format(j))
		dft = df.transpose()
		sums = dft[dft.columns[-1]]
		sums.to_csv('sums_test.csv')
		stats = sums / sums.sum() * 100
		ax = stats.plot.barh(title='Affected land cover (%)', figsize=(22,4))
		ax.figure.savefig('{}.png'.format(j))
		plt.clf()

 # PM input must be the ba_oracle_export shapefile (BA_202X that I load on ftp)


def main():
	# uncomment the following and comment the previous in case the database connection has problems
	# df_sql = pd.read_csv('C:\\Users\piermaio\Documents\gisdata\jrc\weekly_videoconference\df_sql.csv')
	path = 'C:\\Users\piermaio\Documents\gisdata\jrc\weekly_videoconference\\' #PM local output path
	df_sql, nat2k_year, nat2k_week, gdf_sql, tab_nations = db_connection()
	current_year = dt.datetime.today().year
	years = np.arange(2008,2021)  # right limit excluded
	range = dates_range_set(years)
	ppt_generator(path, years, range, df_sql, nat2k_year, nat2k_week)
	# weekly_sum_count(years, range, path, df_sql)
	# PM uncomment this part to generate the corine statistics
	# ba_path = 'C:\\Users\\piermaio\\Documents\\gisdata\\jrc\\BA_2020.shp'  # PM for test
	corine_path = 'C:\\Users\\piermaio\\Documents\\gisdata\\jrc\\BAmapping\\Corine\\raster\\Corine_globcover_MA_TU_ukraine.tif'  # PM local path to the corine geotif
	dict_path = 'C:\\Users\\piermaio\\Documents\\gisdata\\jrc\\weekly_videoconference' # PM local path to the corine key
	corine_stats(gdf_sql, corine_path, dict_path, tab_nations)


if __name__ == "__main__":
	main()
